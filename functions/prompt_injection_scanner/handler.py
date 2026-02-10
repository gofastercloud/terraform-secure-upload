"""
Lambda function that scores text content for prompt injection attacks using
the protectai/deberta-v3-base-prompt-injection-v2 ONNX model.

Input:  {"bucket": "ingress-bucket", "key": "uploads/dave/report.docx"}
Output: {"score": 85, "scannable": true, "format": "docx", "chunks_scanned": 3}
"""

import json
import logging
import os
import tempfile

import boto3
import numpy as np
import onnxruntime as ort
from transformers import AutoTokenizer

logger = logging.getLogger()
logger.setLevel(logging.INFO)

s3 = boto3.client("s3")

MODEL_PATH = os.environ.get("MODEL_PATH", "/opt/model")
MAX_TOKENS = 512
STRIDE = 256

# File extensions we can extract text from
TEXT_EXTENSIONS = {
    ".txt", ".csv", ".md", ".json", ".xml", ".html",
    ".yaml", ".yml", ".log",
}
DOC_EXTENSIONS = {".pdf", ".docx", ".pptx"}
SUPPORTED_EXTENSIONS = TEXT_EXTENSIONS | DOC_EXTENSIONS

# Lazy-loaded globals for model and tokenizer (persist across warm invocations)
_tokenizer = None
_session = None


def _load_model():
    """Load the ONNX model and tokenizer on first invocation."""
    global _tokenizer, _session
    if _tokenizer is None:
        logger.info("Loading tokenizer from %s", MODEL_PATH)
        _tokenizer = AutoTokenizer.from_pretrained(MODEL_PATH)
    if _session is None:
        onnx_path = os.path.join(MODEL_PATH, "onnx", "model.onnx")
        logger.info("Loading ONNX model from %s", onnx_path)
        _session = ort.InferenceSession(
            onnx_path,
            providers=["CPUExecutionProvider"],
        )
    return _tokenizer, _session


def _get_extension(key):
    """Extract lowercase file extension from an S3 object key."""
    _, _, ext = key.rpartition(".")
    return f".{ext.lower()}" if ext else ""


def _extract_text_plain(file_path):
    """Read a plain text file as UTF-8."""
    with open(file_path, "r", encoding="utf-8", errors="replace") as f:
        return f.read()


def _extract_text_pdf(file_path):
    """Extract text from a PDF using PyPDF2."""
    from PyPDF2 import PdfReader

    reader = PdfReader(file_path)
    pages = []
    for page in reader.pages:
        text = page.extract_text()
        if text:
            pages.append(text)
    return "\n".join(pages)


def _extract_text_docx(file_path):
    """Extract text from a .docx file."""
    from docx import Document

    doc = Document(file_path)
    return "\n".join(para.text for para in doc.paragraphs if para.text)


def _extract_text_pptx(file_path):
    """Extract text from a .pptx file."""
    from pptx import Presentation

    prs = Presentation(file_path)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        texts.append(text)
    return "\n".join(texts)


def _extract_text(file_path, ext):
    """Route to the appropriate text extraction method."""
    if ext in TEXT_EXTENSIONS:
        return _extract_text_plain(file_path)
    elif ext == ".pdf":
        return _extract_text_pdf(file_path)
    elif ext == ".docx":
        return _extract_text_docx(file_path)
    elif ext == ".pptx":
        return _extract_text_pptx(file_path)
    return ""


def _score_text(text, tokenizer, session):
    """
    Tokenize text into overlapping chunks and run ONNX inference.
    Returns (max_score, chunks_scanned) where score is 0-100.
    """
    encoding = tokenizer(
        text,
        truncation=False,
        return_attention_mask=False,
        add_special_tokens=False,
    )
    all_ids = encoding["input_ids"]

    if not all_ids:
        return 0, 0

    # Split into overlapping chunks of MAX_TOKENS with STRIDE overlap.
    # Each chunk gets [CLS] ... [SEP] added by the model's tokenizer.
    chunks = []
    start = 0
    while start < len(all_ids):
        end = min(start + MAX_TOKENS - 2, len(all_ids))  # -2 for CLS/SEP
        chunks.append(all_ids[start:end])
        if end >= len(all_ids):
            break
        start += STRIDE

    max_score = 0.0
    for chunk_ids in chunks:
        # Re-encode with special tokens for this chunk
        inputs = tokenizer.encode_plus(
            tokenizer.decode(chunk_ids),
            truncation=True,
            max_length=MAX_TOKENS,
            padding="max_length",
            return_tensors="np",
        )

        ort_inputs = {
            "input_ids": inputs["input_ids"].astype(np.int64),
            "attention_mask": inputs["attention_mask"].astype(np.int64),
        }

        # Check if model expects token_type_ids
        input_names = [inp.name for inp in session.get_inputs()]
        if "token_type_ids" in input_names:
            ort_inputs["token_type_ids"] = np.zeros_like(
                inputs["input_ids"], dtype=np.int64
            )

        outputs = session.run(None, ort_inputs)
        logits = outputs[0]

        # Apply softmax to get probabilities
        exp_logits = np.exp(logits - np.max(logits, axis=-1, keepdims=True))
        probs = exp_logits / np.sum(exp_logits, axis=-1, keepdims=True)

        # INJECTION class is index 1 for this model
        injection_prob = float(probs[0][1])
        score = injection_prob * 100

        if score > max_score:
            max_score = score

    return round(max_score, 1), len(chunks)


def handler(event, context):
    """Lambda entry point."""
    bucket = event["bucket"]
    key = event["key"]

    logger.info("Scanning s3://%s/%s for prompt injection", bucket, key)

    ext = _get_extension(key)

    if ext not in SUPPORTED_EXTENSIONS:
        logger.info(
            "File extension '%s' is not scannable. Returning score 0.", ext
        )
        return {
            "score": 0,
            "scannable": False,
            "format": ext.lstrip(".") if ext else "unknown",
            "chunks_scanned": 0,
        }

    # Download file to /tmp
    with tempfile.NamedTemporaryFile(
        suffix=ext, delete=False, dir="/tmp"
    ) as tmp:
        tmp_path = tmp.name

    try:
        s3.download_file(bucket, key, tmp_path)
        logger.info("Downloaded s3://%s/%s to %s", bucket, key, tmp_path)

        text = _extract_text(tmp_path, ext)

        if not text or not text.strip():
            logger.info("No text extracted from %s. Returning score 0.", key)
            return {
                "score": 0,
                "scannable": False,
                "format": ext.lstrip("."),
                "chunks_scanned": 0,
            }

        tokenizer, session = _load_model()
        score, chunks_scanned = _score_text(text, tokenizer, session)

        logger.info(
            "Scan result for %s: score=%s, chunks_scanned=%d",
            key,
            score,
            chunks_scanned,
        )

        return {
            "score": score,
            "scannable": True,
            "format": ext.lstrip("."),
            "chunks_scanned": chunks_scanned,
        }
    except Exception:
        logger.exception("Error scanning %s", key)
        raise
    finally:
        try:
            os.unlink(tmp_path)
        except OSError:
            pass
